import os
from docx import Document
import json
from PyPDF2 import PdfReader
from pptx import Presentation


def check_file_contents(filename, search_string, case_sensitive=True, exact_match=False):
    """
    Check the contents of a file for the presence of a search string.

    Args:
        filename (str): The name of the file to search.
        search_string (str): The string to search for in the file.
        case_sensitive (bool, optional): Whether the search should be case-sensitive. Defaults to True.
        exact_match (bool, optional): Whether to search for an exact match of the search string. Defaults to False.
    """
    file_extension = os.path.splitext(filename)[1]
    script_path = os.path.abspath(__file__)

    original_search_string = search_string  # Store the original search string

    if not case_sensitive:
        search_string = search_string.lower()

    if file_extension == '.docx':
        doc = Document(filename)
        for paragraph in doc.paragraphs:
            if not case_sensitive:
                paragraph_text = paragraph.text.lower()
            else:
                paragraph_text = paragraph.text

            if exact_match:
                if search_string == paragraph_text:
                    print(f"Found '{original_search_string}' in {filename} - Line: {paragraph.text}")
            else:
                if search_string in paragraph_text:
                    print(f"Found '{original_search_string}' in {filename} - Line: {paragraph.text}")

    elif filename == script_path:
        return

    elif file_extension == '.py':
        with open(filename, 'r') as file:
            for line_number, line in enumerate(file, start=1):
                if not case_sensitive:
                    line = line.lower()

                if exact_match:
                    if search_string == line.strip():
                        print(f"Found '{original_search_string}' in {filename} - Line {line_number}: {line.strip()}")
                else:
                    if search_string in line:
                        print(f"Found '{original_search_string}' in {filename} - Line {line_number}: {line.strip()}")

    elif file_extension == '.txt':
        with open(filename, 'r') as file:
            for line_number, line in enumerate(file, start=1):
                if not case_sensitive:
                    line = line.lower()

                if exact_match:
                    if search_string == line.strip():
                        print(f"Found '{original_search_string}' in {filename} - Line {line_number}: {line.strip()}")
                else:
                    if search_string in line:
                        print(f"Found '{original_search_string}' in {filename} - Line {line_number}: {line.strip()}")

    elif file_extension == '.json':
        with open(filename, 'r') as file:
            data = json.load(file)
            if not case_sensitive:
                data_str = str(data).lower()
            else:
                data_str = str(data)

            if exact_match:
                if search_string == data_str:
                    print(f"Found '{original_search_string}' in {filename}")
            else:
                if search_string in data_str:
                    print(f"Found '{original_search_string}' in {filename}")

    elif file_extension == '.pdf':
        with open(filename, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page_number, page in enumerate(pdf_reader.pages, start=1):
                page_content = page.extract_text()
                if not case_sensitive:
                    page_content = page_content.lower()

                if exact_match:
                    if search_string == page_content.strip():
                        print(f"Found '{original_search_string}' in {filename} - Page {page_number}")
                else:
                    if search_string in page_content:
                        print(f"Found '{original_search_string}' in {filename} - Page {page_number}")

    elif file_extension == '.pptx':
        prs = Presentation(filename)
        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content.append(shape.text)
            slide_text = '\n'.join(slide_content)

            if not case_sensitive:
                slide_text = slide_text.lower()

            if exact_match:
                if search_string == slide_text.strip():
                    print(f"Found '{original_search_string}' in {filename} - Slide {slide_number}")
            else:
                if search_string in slide_text:
                    print(f"Found '{original_search_string}' in {filename} - Slide {slide_number}")

    else:
        print(f"Unsupported file format: {file_extension}")


def search_in_files(search_string):
    """
    Search for a given string in all files within the specified directory or the current directory.

    Args:
        search_string (str): The string to search for.
    """
    case_sensitive_input = input("Do you want the search to be case-sensitive? (y/n) ")
    case_sensitive = case_sensitive_input.lower() == 'y'

    exact_match_input = input("Do you want to search for an exact match? (y/n) ")
    exact_match = exact_match_input.lower() == 'y'

    # Get the directory to search
    search_directory = input("Enter the directory to search (leave blank for current directory): ").strip()
    if not search_directory:
        search_directory = os.path.dirname(os.path.abspath(__file__))

    # Iterate over all files in the directory
    for filename in os.listdir(search_directory):
        # Check if the file is a regular file and exclude temporary files
        if os.path.isfile(os.path.join(search_directory, filename)) and not filename.startswith('~$'):
            check_file_contents(os.path.join(search_directory, filename), search_string, case_sensitive, exact_match)

if __name__ == '__main__':
    """
    Main entry point of the program. Prompts the user for a search string and calls the search_in_files method.
    """
    search_string = input("What string do you want to search for? ")
    search_in_files(search_string)
