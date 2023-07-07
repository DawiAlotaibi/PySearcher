import os
import tkinter as tk
from tkinter import filedialog
from tkinter import ttk
from tkinter import messagebox
from docx import Document
import json
from PyPDF2 import PdfReader
from pptx import Presentation


class App(ttk.Frame):
    def __init__(self, parent):
        ttk.Frame.__init__(self)
        self.treeview_data = []
        self.popup = None
        self.directory_label = None
        self.directory_frame = None
        self.directory_entry = None
        self.browse_button = None
        self.search_button = None
        self.results_label = None
        self.results_text = None
        self.directory_var = None
        self.exact_match_checkbox = None
        self.exact_match_var = None
        self.case_sensitive_checkbox = None
        self.case_sensitive_var = None
        self.search_entry = None
        self.search_label = None
        self.setup_widgets()

    def setup_widgets(self):
        # self.search_label = tk.Label(window, text="Search String:")
        # self.search_label.pack()
        self.search_label = ttk.Label(
            self,
            text="Search String:",
            justify="center",
            font=("-size", 10, "-weight", "bold"),
        )
        self.search_label.grid(row=0, column=0, padx=5, pady=(0, 10), sticky="ew")

        # self.search_entry = tk.Entry(window, width=50)
        # self.search_entry.pack()
        self.search_entry = ttk.Entry(self)
        self.search_entry.insert(0, "Entry")
        self.search_entry.grid(row=0, column=1, padx=5, pady=(0, 10), sticky="ew")

        self.case_sensitive_var = tk.BooleanVar()
        # self.case_sensitive_checkbox = tk.Checkbutton(window, text="Case Sensitive", variable=self.case_sensitive_var)
        # self.case_sensitive_checkbox.pack()
        self.case_sensitive_checkbox = ttk.Checkbutton(
            self, text="Case Sensitive", variable=self.case_sensitive_var
        )
        self.case_sensitive_checkbox.grid(row=0, column=2, padx=5, pady=(0, 10), sticky="ew")

        self.exact_match_var = tk.BooleanVar()
        # self.exact_match_checkbox = tk.Checkbutton(window, text="Exact Match", variable=self.exact_match_var)
        # self.exact_match_checkbox.pack()
        self.exact_match_checkbox = ttk.Checkbutton(
            self, text="Exact Match", variable=self.exact_match_var
        )
        self.exact_match_checkbox.grid(row=0, column=3, padx=5, pady=(0, 10), sticky="ew")

        self.directory_var = tk.StringVar()
        # self.directory_label = tk.Label(window, text="Search Directory:")
        # self.directory_label.pack()
        self.directory_label = ttk.Label(
            self,
            text="Search Directory:",
            justify="center",
            font=("-size", 10, "-weight", "bold"),
        )
        self.directory_label.grid(row=1, column=0, padx=5, pady=(0, 10), sticky="ew")

        # self.directory_frame = tk.Frame(window)
        # self.directory_entry = tk.Entry(self.directory_frame, textvariable=self.directory_var, width=40)
        # self.directory_entry.pack(side=tk.LEFT)
        # self.directory_frame = ttk.Frame(self, padding=(0, 0, 0, 10))
        # self.directory_frame.grid(row=1, column=1, padx=5, pady=(0, 10), sticky="ew", rowspan=3)
        self.directory_entry = ttk.Entry(self, textvariable=self.directory_var, width=40)
        self.directory_entry.grid(row=1, column=1, padx=5, pady=(0, 10), sticky="ew")
        self.browse_button = ttk.Button(self, text="Browse", command=browse_directory)
        self.browse_button.grid(row=1, column=2, padx=5, pady=(0, 10), sticky="ew")
        self.search_button = ttk.Button(self, text="Search", command=search_files)
        self.search_button.grid(row=2, column=1, padx=5, pady=(30, 10), sticky="ew")

        # self.browse_button = tk.Button(self.directory_frame, text="Browse", command=browse_directory)
        # self.browse_button.pack(side=tk.LEFT)
        #
        # self.directory_frame.pack()
        #
        # self.search_button = tk.Button(window, text="Search", command=search_files)
        # self.search_button.pack()
        #
        # self.results_label = tk.Label(window, text="Results:")
        # self.results_label.pack()
        #
        self.results_text = tk.Text(window, width=60, height=15)
        # self.results_text.pack()

    def clear_treeview(self):
        self.popup.treeview.delete(*self.popup.treeview.get_children())
        self.treeview_data.clear()
        self.popup.destroy()

    def show_popup(self, original_search_string):
        self.popup = tk.Toplevel()
        self.popup.title("Results")
        self.popup.protocol("WM_DELETE_WINDOW", self.clear_treeview)
        self.popup.searchString = ttk.Label(
            self,
            text=original_search_string,
            justify="center",
            font=("-size", 10, "-weight", "bold"),
        )
        self.search_label.grid(row=0, column=0, padx=5, pady=(0, 10), sticky="ew")

        self.popup.scrollbar = ttk.Scrollbar(self.popup)
        self.popup.scrollbar.pack(side="right", fill="y")
        self.popup.treeview = ttk.Treeview(
            self.popup,
            selectmode="browse",
            yscrollcommand=self.popup.scrollbar.set,
            columns=(1, 2),
            height=10,
        )
        self.popup.treeview.pack(expand=True, fill="both")
        self.popup.scrollbar.config(command=self.popup.treeview.yview)

        # Treeview columns
        self.popup.treeview.column("#0", anchor="w", width=120)
        self.popup.treeview.column(2, anchor="w", width=120)
        self.popup.treeview.column(1, anchor="w", width=120)

        # Treeview headings
        self.popup.treeview.heading("#0", text="Path", anchor="center")
        self.popup.treeview.heading(1, text="File", anchor="center")
        self.popup.treeview.heading(2, text="Line/Page", anchor="center")
        # # Add content to the popup window
        # label = ttk.Label(self.popup, text="This is a popup window!")
        # label.pack(padx=10, pady=10)
        #
        # # Add a button to close the popup
        # close_button = ttk.Button(self.popup, text="Close", command=self.popup.destroy)
        # close_button.pack(pady=10,fill="both", expand=True)

        for item in self.treeview_data:
            print(item)
            self.popup.treeview.insert(
                parent=item[0], index="end", iid=item[1], text=item[2], values=item[3]
            )
            if item[0] == "" or item[1] in {8, 21}:
                self.popup.treeview.item(item[1], open=True)  # Open parents
        self.popup.update()
        self.popup.minsize(self.popup.winfo_width(), self.popup.winfo_height())
        x_coordinate = int((self.popup.winfo_screenwidth() / 2) - (self.popup.winfo_width() / 2))
        y_coordinate = int((self.popup.winfo_screenheight() / 2) - (self.popup.winfo_height() / 2))
        self.popup.geometry("+{}+{}".format(x_coordinate, y_coordinate - 20))


def check_file_contents(filename, search_string, case_sensitive=True, exact_match=False):
    file_extension = os.path.splitext(filename)[1]
    script_path = os.path.abspath(__file__)
    original_search_string = search_string

    if not case_sensitive:
        search_string = search_string.lower()

    if file_extension == '.docx':
        doc = Document(filename)
        line_number = 1
        for paragraph in doc.paragraphs:
            if not case_sensitive:
                paragraph_text = paragraph.text.lower()
            else:
                paragraph_text = paragraph.text

            if exact_match:
                if search_string == paragraph_text:
                    # app.results_text.insert(tk.END,
                    #                         f"Found '{original_search_string}' in {filename} - Line: {line_number}\n")
                    actualFileName = os.path.basename(filename)
                    app.treeview_data.append(
                        ("", len(app.treeview_data) + 1, filename, (actualFileName, f"Line: {line_number}")))
            else:
                if search_string in paragraph_text:
                    # app.results_text.insert(tk.END,
                    #                         f"Found '{original_search_string}' in {filename} - Line: {line_number}\n")
                    actualFileName = os.path.basename(filename)
                    app.treeview_data.append(
                        ("", len(app.treeview_data) + 1, filename, (actualFileName, f"Line: {line_number}")))

            line_number += 1


    elif file_extension == '.py':
        with open(filename, 'r') as file:
            for line_number, line in enumerate(file, start=1):
                if not case_sensitive:
                    line = line.lower()

                if exact_match:
                    if search_string == line.strip():
                        # app.results_text.insert(tk.END,
                        #                         f"Found '{original_search_string}' in {filename} - Line {line_number}: {line.strip()}\n")
                        actualFileName = os.path.basename(filename)
                        app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                                  (actualFileName, f"Line: {line_number}")))

                else:
                    if search_string in line:
                        # app.results_text.insert(tk.END, f"Found '{original_search_string}' in {filename} - Line {
                        # line_number}: {line.strip()}\n")
                        actualFileName = os.path.basename(filename)
                        app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                                  (actualFileName, f"Line: {line_number}")))



    elif file_extension == '.txt':
        with open(filename, 'r') as file:
            for line_number, line in enumerate(file, start=1):
                if not case_sensitive:
                    line = line.lower()

                if exact_match:
                    if search_string == line.strip():
                        # app.results_text.insert(tk.END,
                        #                         f"Found '{original_search_string}' in {filename} - Line {line_number}: {line.strip()}\n")
                        actualFileName = os.path.basename(filename)
                        app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                                  (actualFileName, f"Line: {line_number}")))
                else:
                    if search_string in line:
                        # app.results_text.insert(tk.END,
                        #                         f"Found '{original_search_string}' in {filename} - Line {line_number}: {line.strip()}\n")
                        actualFileName = os.path.basename(filename)
                        app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                                  (actualFileName, f"Line: {line_number}")))

    elif file_extension == '.json':  # I need to find a way to print the line here
        with open(filename, 'r') as file:
            data = json.load(file)
            if not case_sensitive:
                data_str = str(data).lower()
            else:
                data_str = str(data)

            if exact_match:
                if search_string == data_str:
                    # app.results_text.insert(tk.END, f"Found '{original_search_string}' in {filename}\n")
                    actualFileName = os.path.basename(filename)
                    app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                              (actualFileName, data_str)))
            else:
                if search_string in data_str:
                    # app.results_text.insert(tk.END, f"Found '{original_search_string}' in {filename}\n")
                    actualFileName = os.path.basename(filename)
                    app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                              (actualFileName, data_str)))

    elif file_extension == '.pdf':
        with open(filename, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page_number, page in enumerate(pdf_reader.pages, start=1):
                page_content = page.extract_text()
                if not case_sensitive:
                    page_content = page_content.lower()

                if exact_match:
                    if search_string == page_content.strip():
                        actualFileName = os.path.basename(filename)
                        app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                                  (actualFileName, f"Page: {page_number}")))

                else:
                    if search_string in page_content:
                        actualFileName = os.path.basename(filename)
                        app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                                  (actualFileName, f"Page: {page_number}")))
    elif file_extension == '.pptx':
        prs = Presentation(filename)
        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content.append(shape.text)
            slide_text = '\n'.join(slide_content)

            if not case_sensitive:
                slide_text = slide_text.lower()

            if exact_match:
                if search_string == slide_text.strip():
                    actualFileName = os.path.basename(filename)
                    app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                              (actualFileName, f"Page: {slide_number}")))
            else:
                if search_string in slide_text:
                    actualFileName = os.path.basename(filename)
                    app.treeview_data.append(("", len(app.treeview_data) + 1, filename,
                                              (actualFileName, f"Page: {slide_number}")))
    else:
        app.results_text.insert(tk.END, f"Unsupported file format: {file_extension}\n")


def browse_directory():
    selected_directory = filedialog.askdirectory()
    app.directory_var.set(selected_directory)


def search_files():
    search_string = app.search_entry.get().strip()
    if not search_string:
        messagebox.showwarning("Empty Search String", "Please enter a search string.")
        return

    case_sensitive = app.case_sensitive_var.get()
    exact_match = app.exact_match_var.get()

    search_directory = app.directory_var.get()
    if not search_directory:
        messagebox.showwarning("No Directory Selected", "Please select a directory to search.")
        return

    script_file = os.path.normpath(os.path.abspath(__file__))

    for filename in os.listdir(search_directory):
        file_path = os.path.normpath(os.path.join(search_directory, filename))

        if os.path.isfile(file_path) and not filename.startswith('~$') and file_path != script_file:
            check_file_contents(file_path, search_string, case_sensitive, exact_match)
    app.show_popup(search_string)


# Create the main window
window = tk.Tk()
window.title("File Searcher")

window.tk.call("source", "azure.tcl")
window.tk.call("set_theme", "dark")

app = App(window)
app.pack(fill="both", expand=True)

window.update()
window.minsize(window.winfo_width(), window.winfo_height())
x_coordinate = int((window.winfo_screenwidth() / 2) - (window.winfo_width() / 2))
y_coordinate = int((window.winfo_screenheight() / 2) - (window.winfo_height() / 2))
window.geometry("+{}+{}".format(x_coordinate, y_coordinate - 20))

window.mainloop()
